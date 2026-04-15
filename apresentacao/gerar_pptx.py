from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD = RGBColor(0x25, 0x25, 0x3A)
ACCENT_BLUE = RGBColor(0x4F, 0xC3, 0xF7)
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)
ACCENT_PURPLE = RGBColor(0xAB, 0x47, 0xBC)
ACCENT_GREEN = RGBColor(0x66, 0xBB, 0x6A)
ACCENT_RED = RGBColor(0xEF, 0x53, 0x50)
ACCENT_YELLOW = RGBColor(0xFD, 0xD8, 0x35)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBD, 0xBD, 0xBD)
DARK_TEXT = RGBColor(0x21, 0x21, 0x21)


def set_slide_bg(slide, color):
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = color
    bg_shape.line.fill.background()
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, w, h, fill_color, text="", font_size=12, font_color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = align
    return shape


def add_rect(slide, left, top, w, h, fill_color, text="", font_size=11, font_color=WHITE, bold=False, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        tf.margin_left = Pt(4)
        tf.margin_right = Pt(4)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = align
    return shape


def add_text(slide, left, top, w, h, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_bullet_list(slide, left, top, w, h, items, font_size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return txBox


def add_arrow_down(slide, x, y, length, color=LIGHT_GRAY):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.03), length)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, x - Inches(0.08), y + length, Inches(0.2), Inches(0.15))
    tri.fill.solid()
    tri.fill.fore_color.rgb = color
    tri.line.fill.background()
    tri.rotation = 180.0


def add_arrow_right(slide, x, y, length, color=LIGHT_GRAY):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, length, Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, x + length, y - Inches(0.08), Inches(0.15), Inches(0.2))
    tri.fill.solid()
    tri.fill.fore_color.rgb = color
    tri.line.fill.background()
    tri.rotation = 90.0


# ============================================================
# SLIDE 1 — COVER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BLUE)
add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         "Intelligent Streaming Channel\nManagement Platform with AI",
         font_size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
         "AI-Powered Chatbot for NOC Operators", font_size=24, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
         "Amazon Bedrock  |  DynamoDB  |  MediaLive  |  MediaPackage V2  |  MediaTailor  |  CloudFront",
         font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08), ACCENT_PURPLE)


# ============================================================
# SLIDE 2 — THE CHALLENGE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_RED)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
         "The Challenge", font_size=32, color=WHITE, bold=True)

problems = [
    ("~220 live channels", "NOC operators manage hundreds of live streaming channels simultaneously"),
    ("Multiple AWS consoles", "Querying configurations, logs and metrics requires navigating multiple consoles"),
    ("Complex channel creation", "Creating a channel involves 4+ AWS services with interdependent configurations"),
    ("Manual troubleshooting", "Manual correlation of data from different sources to diagnose problems"),
]
for i, (title, desc) in enumerate(problems):
    y = Inches(1.5) + Inches(i * 1.35)
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(1.1), BG_CARD)
    add_rect(slide, Inches(0.8), y, Inches(0.08), Inches(1.1), ACCENT_RED)
    add_text(slide, Inches(1.2), y + Inches(0.1), Inches(10), Inches(0.4),
             title, font_size=20, color=ACCENT_RED, bold=True)
    add_text(slide, Inches(1.2), y + Inches(0.55), Inches(10), Inches(0.4),
             desc, font_size=16, color=LIGHT_GRAY)


# ============================================================
# SLIDE 3 — THE SOLUTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_GREEN)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
         "The Solution — Intelligent NOC Chatbot", font_size=30, color=WHITE, bold=True)

solutions = [
    ("Natural Language", "Conversational interface in Portuguese for queries and operations"),
    ("Intelligent RAG", "Retrieval-Augmented Generation over real channel configs and metrics"),
    ("DynamoDB Fast Queries", "Sub-second queries via DynamoDB with S3 fallback (was 30-120s, now <500ms)"),
    ("Automated Creation", "Orchestrated channel creation with template cloning and automatic rollback"),
    ("Proactive Alerts", "SNS notifications for ERROR/CRITICAL events with suppression window"),
    ("Data Export", "CSV/JSON generation with pre-signed download URLs"),
]
for i, (title, desc) in enumerate(solutions):
    y = Inches(1.4) + Inches(i * 1.1)
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.9), BG_CARD)
    add_rect(slide, Inches(0.8), y, Inches(0.08), Inches(0.9), ACCENT_GREEN)
    add_text(slide, Inches(1.2), y + Inches(0.08), Inches(3), Inches(0.35),
             title, font_size=18, color=ACCENT_GREEN, bold=True)
    add_text(slide, Inches(1.2), y + Inches(0.45), Inches(10), Inches(0.35),
             desc, font_size=15, color=LIGHT_GRAY)


# ============================================================
# SLIDE 4 — DIAGRAM: OVERALL ARCHITECTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BLUE)
add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
         "Solution Architecture", font_size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# User
add_shape(slide, Inches(5.8), Inches(0.55), Inches(1.8), Inches(0.5), RGBColor(0x37,0x47,0x4F),
          "NOC Operator", font_size=11, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_arrow_down(slide, Inches(6.65), Inches(1.05), Inches(0.2), LIGHT_GRAY)

# Presentation Layer
LAYER_BLUE = RGBColor(0x1A, 0x3A, 0x5C)
add_rect(slide, Inches(0.3), Inches(1.3), Inches(12.7), Inches(0.85), LAYER_BLUE)
add_text(slide, Inches(0.4), Inches(1.3), Inches(2), Inches(0.3),
         "Presentation", font_size=9, color=ACCENT_BLUE, bold=True)
add_shape(slide, Inches(1.5), Inches(1.45), Inches(1.8), Inches(0.55), BG_CARD,
          "CloudFront\n(CDN)", font_size=10, font_color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_shape(slide, Inches(4.0), Inches(1.45), Inches(1.8), Inches(0.55), BG_CARD,
          "S3 Frontend\n(HTML/JS)", font_size=10, font_color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_shape(slide, Inches(6.5), Inches(1.45), Inches(1.8), Inches(0.55), BG_CARD,
          "Cognito\n(JWT Auth)", font_size=10, font_color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(9.0), Inches(1.5), Inches(3.5), Inches(0.4),
         "Dark theme chat UI, 45+ categorized suggestions", font_size=9, color=LIGHT_GRAY)
add_arrow_down(slide, Inches(6.65), Inches(2.15), Inches(0.15), LIGHT_GRAY)

# API Layer
LAYER_ORANGE = RGBColor(0x4A, 0x2E, 0x0A)
add_rect(slide, Inches(0.3), Inches(2.35), Inches(12.7), Inches(0.85), LAYER_ORANGE)
add_text(slide, Inches(0.4), Inches(2.35), Inches(2), Inches(0.3),
         "API", font_size=9, color=ACCENT_ORANGE, bold=True)
add_shape(slide, Inches(3.5), Inches(2.5), Inches(2.2), Inches(0.55), BG_CARD,
          "API Gateway REST", font_size=10, font_color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_shape(slide, Inches(7.0), Inches(2.5), Inches(2.5), Inches(0.55), BG_CARD,
          "Lambda Function URL\n(5min timeout)", font_size=10, font_color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_arrow_down(slide, Inches(6.65), Inches(3.2), Inches(0.15), LIGHT_GRAY)

# Orchestration Layer
LAYER_PURPLE = RGBColor(0x2E, 0x1A, 0x47)
add_rect(slide, Inches(0.3), Inches(3.4), Inches(12.7), Inches(1.85), LAYER_PURPLE)
add_text(slide, Inches(0.4), Inches(3.4), Inches(3), Inches(0.3),
         "Orchestration (us-east-1)", font_size=9, color=ACCENT_PURPLE, bold=True)
add_shape(slide, Inches(5.4), Inches(3.6), Inches(2.5), Inches(0.55), BG_CARD,
          "Orchestrator Lambda\n(Entry Point)", font_size=10, font_color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_arrow_down(slide, Inches(3.2), Inches(4.15), Inches(0.15), ACCENT_PURPLE)
add_arrow_down(slide, Inches(6.65), Inches(4.15), Inches(0.15), ACCENT_GREEN)
add_arrow_down(slide, Inches(10.2), Inches(4.15), Inches(0.15), ACCENT_ORANGE)

add_shape(slide, Inches(1.5), Inches(4.35), Inches(2.0), Inches(0.55), BG_CARD,
          "Bedrock Agent\n(Claude PT-BR)", font_size=10, font_color=ACCENT_PURPLE, bold=True, align=PP_ALIGN.CENTER)
add_shape(slide, Inches(3.8), Inches(4.35), Inches(1.5), Inches(0.55), BG_CARD,
          "KB_CONFIG\n(~220 channels)", font_size=9, font_color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
add_shape(slide, Inches(3.8), Inches(4.95), Inches(1.5), Inches(0.45), BG_CARD,
          "KB_LOGS\n(metrics)", font_size=9, font_color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.8), Inches(4.35), Inches(2.0), Inches(0.55), BG_CARD,
          "Exporter Lambda\n(CSV/JSON)", font_size=10, font_color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)
add_shape(slide, Inches(9.0), Inches(4.35), Inches(2.5), Inches(0.55), BG_CARD,
          "Configurator Lambda\n(Create Channels + Rollback)", font_size=10, font_color=ACCENT_RED, bold=True, align=PP_ALIGN.CENTER)
add_arrow_right(slide, Inches(3.5), Inches(4.6), Inches(0.25), ACCENT_GREEN)
add_arrow_down(slide, Inches(6.65), Inches(5.25), Inches(0.15), LIGHT_GRAY)

# Storage Layer
LAYER_GREEN = RGBColor(0x1A, 0x3A, 0x1A)
add_rect(slide, Inches(0.3), Inches(5.45), Inches(12.7), Inches(0.7), LAYER_GREEN)
add_text(slide, Inches(0.4), Inches(5.45), Inches(2), Inches(0.3),
         "Storage", font_size=9, color=ACCENT_GREEN, bold=True)
buckets = [
    ("S3_KBConfig\n(flat JSON)", Inches(0.6)),
    ("S3_KBLogs\n(Events)", Inches(2.5)),
    ("DynamoDB\nStreamingConfigs", Inches(4.4)),
    ("DynamoDB\nStreamingLogs", Inches(6.3)),
    ("S3_Audit\n(365 days)", Inches(8.2)),
    ("S3_Exports\n(24h TTL)", Inches(10.2)),
]
for name, x in buckets:
    color = RGBColor(0x0D, 0x47, 0xA1) if "DynamoDB" in name else BG_CARD
    fc = ACCENT_BLUE if "DynamoDB" in name else ACCENT_GREEN
    add_shape(slide, x, Inches(5.55), Inches(1.7), Inches(0.5), color,
              name, font_size=9, font_color=fc, align=PP_ALIGN.CENTER)
add_arrow_down(slide, Inches(6.65), Inches(6.15), Inches(0.1), LIGHT_GRAY)

# Ingestion + Streaming
LAYER_YELLOW = RGBColor(0x3A, 0x35, 0x0A)
add_rect(slide, Inches(0.3), Inches(6.3), Inches(5.8), Inches(1.05), LAYER_YELLOW)
add_text(slide, Inches(0.4), Inches(6.3), Inches(2), Inches(0.3),
         "Ingestion", font_size=9, color=ACCENT_YELLOW, bold=True)
add_shape(slide, Inches(0.5), Inches(6.55), Inches(2.3), Inches(0.4), BG_CARD,
          "EventBridge (6h) > Pipeline_Config", font_size=9, font_color=ACCENT_YELLOW, align=PP_ALIGN.CENTER)
add_shape(slide, Inches(3.2), Inches(6.55), Inches(2.7), Inches(0.4), BG_CARD,
          "EventBridge (1h) > Pipeline_Logs", font_size=9, font_color=ACCENT_YELLOW, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(6.98), Inches(5.5), Inches(0.3),
         "Dual-write: S3 (RAG) + DynamoDB (fast queries) | Configs 6h, Metrics 1h",
         font_size=8, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

LAYER_RED = RGBColor(0x4A, 0x1A, 0x1A)
add_rect(slide, Inches(6.5), Inches(6.3), Inches(6.5), Inches(1.05), LAYER_RED)
add_text(slide, Inches(6.6), Inches(6.3), Inches(3), Inches(0.3),
         "Streaming (sa-east-1)", font_size=9, color=ACCENT_RED, bold=True)
for name, x in [("MediaLive", Inches(6.7)), ("MediaPkg V2", Inches(8.2)), ("MediaTailor", Inches(9.7)), ("CloudFront", Inches(11.2))]:
    add_shape(slide, x, Inches(6.55), Inches(1.3), Inches(0.4), BG_CARD,
              name, font_size=9, font_color=ACCENT_RED, bold=True, align=PP_ALIGN.CENTER)
add_shape(slide, Inches(6.7), Inches(7.0), Inches(1.3), Inches(0.3), BG_CARD,
          "CloudWatch", font_size=9, font_color=ACCENT_YELLOW, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5 — DIAGRAM: CHATBOT QUERY FLOW
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_PURPLE)
add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
         "Flow Diagram — Chatbot Query (RAG)", font_size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

BOX_H = Inches(0.9)
Y_MAIN = Inches(2.8)

add_shape(slide, Inches(0.3), Y_MAIN, Inches(1.3), BOX_H, RGBColor(0x37,0x47,0x4F),
          "NOC\nOperator", font_size=12, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_arrow_right(slide, Inches(1.6), Y_MAIN + Inches(0.4), Inches(0.4), ACCENT_BLUE)

add_shape(slide, Inches(2.1), Y_MAIN, Inches(1.6), BOX_H, RGBColor(0x0D,0x47,0xA1),
          "CloudFront\n+ S3 Frontend", font_size=11, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_arrow_right(slide, Inches(3.7), Y_MAIN + Inches(0.4), Inches(0.4), ACCENT_ORANGE)

add_shape(slide, Inches(2.1), Inches(1.5), Inches(1.6), Inches(0.6), RGBColor(0x88,0x0E,0x4F),
          "Cognito (JWT)", font_size=10, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_arrow_down(slide, Inches(2.88), Inches(2.1), Inches(0.55), RGBColor(0x88,0x0E,0x4F))

add_shape(slide, Inches(4.2), Y_MAIN, Inches(1.8), BOX_H, RGBColor(0xE6,0x51,0x00),
          "Orchestrator\nLambda\n(5min timeout)", font_size=11, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_arrow_right(slide, Inches(6.0), Y_MAIN + Inches(0.4), Inches(0.4), ACCENT_PURPLE)

add_shape(slide, Inches(6.5), Y_MAIN, Inches(1.8), BOX_H, RGBColor(0x4A,0x14,0x8C),
          "Bedrock Agent\n(Claude)\nPortuguese", font_size=11, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_arrow_right(slide, Inches(8.3), Inches(2.5), Inches(0.5), ACCENT_GREEN)
add_arrow_right(slide, Inches(8.3), Inches(3.8), Inches(0.5), ACCENT_GREEN)

add_shape(slide, Inches(8.9), Inches(1.8), Inches(2.0), Inches(1.1), RGBColor(0x1B,0x5E,0x20),
          "KB_CONFIG\n(~220 channels)\nFlat JSON", font_size=11, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_shape(slide, Inches(8.9), Inches(3.5), Inches(2.0), Inches(1.1), RGBColor(0x1B,0x5E,0x20),
          "KB_LOGS\n(CW Metrics)\nStructured Events", font_size=11, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(8.4), Inches(2.95), Inches(0.5), Inches(0.3),
         "RAG", font_size=10, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(4.2), Inches(5.0), Inches(3.0), Inches(0.8), RGBColor(0x4A,0x14,0x8C),
          "Contextualized response\nin Portuguese", font_size=12, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_arrow_down(slide, Inches(5.1), Inches(3.7), Inches(1.15), ACCENT_PURPLE)
add_rect(slide, Inches(2.8), Inches(5.2), Inches(1.4), Inches(0.03), ACCENT_PURPLE)
add_text(slide, Inches(2.5), Inches(4.85), Inches(1.5), Inches(0.3),
         "<- response", font_size=9, color=ACCENT_PURPLE, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(0.3), Inches(5.0), Inches(2.2), Inches(2.2), BG_CARD)
add_text(slide, Inches(0.4), Inches(5.05), Inches(2), Inches(0.3),
         "Example queries:", font_size=11, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(0.4), Inches(5.35), Inches(2), Inches(1.8),
                ['"Channels with 1080p?"', '"Critical errors 24h?"', '"Config of channel X?"', '"H.265 channels?"', '"Export CSV?"'],
                font_size=9, color=LIGHT_GRAY)

add_text(slide, Inches(11.2), Inches(5.0), Inches(2), Inches(0.3),
         "Legend:", font_size=10, color=WHITE, bold=True)
for i, (c, label) in enumerate([(ACCENT_BLUE, "Presentation"), (ACCENT_ORANGE, "Compute"), (ACCENT_PURPLE, "AI / Bedrock"), (ACCENT_GREEN, "Data / S3")]):
    y = Inches(5.35) + Inches(i * 0.35)
    add_rect(slide, Inches(11.2), y, Inches(0.25), Inches(0.2), c)
    add_text(slide, Inches(11.55), y - Inches(0.02), Inches(1.5), Inches(0.25),
             label, font_size=9, color=LIGHT_GRAY)


# ============================================================
# SLIDE 6 — DIAGRAM: ORCHESTRATED CHANNEL CREATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_ORANGE)
add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
         "Flow Diagram — Orchestrated Channel Creation", font_size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(0.3), Inches(1.0), Inches(1.5), Inches(0.7), RGBColor(0x37,0x47,0x4F),
          "Operator\n\"Create channel X\"", font_size=10, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_arrow_right(slide, Inches(1.8), Inches(1.3), Inches(0.3), LIGHT_GRAY)
add_shape(slide, Inches(2.2), Inches(0.9), Inches(2.3), Inches(0.9), RGBColor(0x4A,0x14,0x8C),
          "Bedrock Agent\nCollects parameters\nconversationally", font_size=10, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_arrow_right(slide, Inches(4.5), Inches(1.3), Inches(0.3), LIGHT_GRAY)
add_shape(slide, Inches(4.9), Inches(0.9), Inches(2.3), Inches(0.9), RGBColor(0xE6,0x51,0x00),
          "Configurator\nLambda", font_size=12, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

STEP_X = Inches(1.5)
STEP_W = Inches(5.5)
STEP_H = Inches(0.7)
steps = [
    ("Step 1", "Create MediaPackage V2 Channel (Channel Group + Channel)", ACCENT_BLUE),
    ("Step 2", "Create HLS / DASH Endpoints on MediaPackage V2", ACCENT_GREEN),
    ("Step 3", "Create MediaLive Inputs (RTMP/RTP with Failover)", ACCENT_ORANGE),
    ("Step 4", "Create MediaLive Channel (links inputs + endpoints)", ACCENT_PURPLE),
]
for i, (step, desc, color) in enumerate(steps):
    y = Inches(2.2) + Inches(i * 1.1)
    add_rect(slide, STEP_X, y, Inches(1.0), STEP_H, color,
             step, font_size=12, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_shape(slide, STEP_X + Inches(1.1), y, STEP_W - Inches(1.1), STEP_H, BG_CARD,
              desc, font_size=12, font_color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    if i < 3:
        add_arrow_down(slide, STEP_X + Inches(0.48), y + STEP_H, Inches(0.25), color)

y_success = Inches(2.2) + Inches(4 * 1.1)
add_shape(slide, STEP_X, y_success, STEP_W, Inches(0.6), RGBColor(0x1B,0x5E,0x20),
          "Channel Created Successfully!", font_size=14, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Rollback
RB_X = Inches(8.0)
add_shape(slide, RB_X, Inches(2.0), Inches(4.8), Inches(3.0), RGBColor(0x4A,0x1A,0x1A))
add_rect(slide, RB_X, Inches(2.0), Inches(4.8), Inches(0.5), ACCENT_RED,
         "AUTOMATIC ROLLBACK", font_size=13, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_bullet_list(slide, RB_X + Inches(0.2), Inches(2.6), Inches(4.4), Inches(2.2),
                ["Failure at any step triggers rollback", "Undoes previous steps in reverse order",
                 "Full audit trail in S3_Audit", "Operator receives detailed error message"],
                font_size=12, color=LIGHT_GRAY)
add_text(slide, Inches(7.1), Inches(3.2), Inches(0.9), Inches(0.4),
         "failure >", font_size=10, color=ACCENT_RED, bold=True, align=PP_ALIGN.CENTER)

# Parameters
add_shape(slide, RB_X, Inches(5.3), Inches(4.8), Inches(2.0), BG_CARD)
add_rect(slide, RB_X, Inches(5.3), Inches(4.8), Inches(0.4), ACCENT_BLUE,
         "Parameters Collected via Chat", font_size=11, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_bullet_list(slide, RB_X + Inches(0.2), Inches(5.8), Inches(4.4), Inches(1.4),
                ["Channel name", "Video codec (H.264 / H.265)", "Resolution (1080p / 720p)",
                 "Bitrate / Input type (RTMP/RTP)", "Automatic failover (yes/no)"],
                font_size=11, color=LIGHT_GRAY)


# ============================================================
# SLIDE 7 — DIAGRAM: DynamoDB MIGRATION (NEW)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
ACCENT_DDB = RGBColor(0x0D, 0x47, 0xA1)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_DDB)
add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
         "DynamoDB — Fast Query Layer", font_size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Before vs After
add_shape(slide, Inches(0.3), Inches(0.85), Inches(6.2), Inches(2.6), BG_CARD)
add_rect(slide, Inches(0.3), Inches(0.85), Inches(6.2), Inches(0.45), ACCENT_RED,
         "BEFORE — S3 Only", font_size=13, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
before_items = [
    "Exporter lists + reads thousands of JSON files from S3",
    "Query latency: 30-120 seconds per request",
    "No indexing — full scan for every query",
    "Chatbot responses slow for data-heavy questions",
]
add_bullet_list(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(1.8),
                before_items, font_size=12, color=LIGHT_GRAY)

add_shape(slide, Inches(6.8), Inches(0.85), Inches(6.2), Inches(2.6), BG_CARD)
add_rect(slide, Inches(6.8), Inches(0.85), Inches(6.2), Inches(0.45), ACCENT_GREEN,
         "AFTER — DynamoDB + S3 Dual-Write", font_size=13, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
after_items = [
    "DynamoDB queries with GSI indexes (<500ms)",
    "S3 preserved for Bedrock RAG Knowledge Base",
    "Automatic fallback to S3 if DynamoDB fails",
    "~$3/month additional cost (PAY_PER_REQUEST)",
]
add_bullet_list(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(1.8),
                after_items, font_size=12, color=LIGHT_GRAY)

# Tables detail
add_shape(slide, Inches(0.3), Inches(3.7), Inches(6.2), Inches(3.5), BG_CARD)
add_rect(slide, Inches(0.3), Inches(3.7), Inches(6.2), Inches(0.45), ACCENT_DDB,
         "StreamingConfigs Table", font_size=13, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
cfg_items = [
    "PK: {servico}#{tipo}  |  SK: {resource_id}",
    "GSI_NomeCanal: search by channel name",
    "~285 items  |  Point-in-time recovery",
    "Top-level: servico, nome_canal, estado, regiao",
    "Full config in 'data' field (JSON)",
]
add_bullet_list(slide, Inches(0.5), Inches(4.25), Inches(5.8), Inches(2.5),
                cfg_items, font_size=11, color=LIGHT_GRAY)

add_shape(slide, Inches(6.8), Inches(3.7), Inches(6.2), Inches(3.5), BG_CARD)
add_rect(slide, Inches(6.8), Inches(3.7), Inches(6.2), Inches(0.45), ACCENT_DDB,
         "StreamingLogs Table", font_size=13, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
log_items = [
    "PK: {servico}#{canal}  |  SK: {timestamp}#{metrica}",
    "GSI_Severidade: filter by severity level",
    "TTL: 30 days auto-expiration",
    "~1.74M items/month  |  ~870MB storage",
    "Top-level: severidade, canal, metrica_nome/valor",
]
add_bullet_list(slide, Inches(7.0), Inches(4.25), Inches(5.8), Inches(2.5),
                log_items, font_size=11, color=LIGHT_GRAY)


# ============================================================
# SLIDE 8 — DIAGRAM: INGESTION PIPELINES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_YELLOW)
add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
         "Flow Diagram — Ingestion Pipelines", font_size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Config Pipeline (left)
add_rect(slide, Inches(0.3), Inches(0.8), Inches(6.2), Inches(0.45), ACCENT_GREEN,
         "Configuration Pipeline (every 6h)", font_size=13, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
cfg_y = Inches(1.5)
add_shape(slide, Inches(0.5), cfg_y, Inches(1.6), Inches(0.65), RGBColor(0x88,0x0E,0x4F),
          "EventBridge\nScheduler (6h)", font_size=10, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_arrow_right(slide, Inches(2.1), cfg_y + Inches(0.3), Inches(0.3), ACCENT_YELLOW)
add_shape(slide, Inches(2.5), cfg_y, Inches(1.6), Inches(0.65), RGBColor(0xE6,0x51,0x00),
          "Pipeline_Config\nLambda", font_size=10, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_arrow_right(slide, Inches(4.1), cfg_y + Inches(0.3), Inches(0.3), ACCENT_YELLOW)
add_shape(slide, Inches(4.5), cfg_y, Inches(1.8), Inches(0.65), RGBColor(0x4A,0x1A,0x1A),
          "AWS APIs\n(4 services)", font_size=10, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, svc in enumerate(["MediaLive", "MediaPackage V2", "MediaTailor", "CloudFront"]):
    add_shape(slide, Inches(0.5) + Inches(i * 1.55), Inches(2.5), Inches(1.4), Inches(0.4), BG_CARD,
              svc, font_size=9, font_color=ACCENT_RED, bold=True, align=PP_ALIGN.CENTER)

add_arrow_down(slide, Inches(3.3), Inches(2.9), Inches(0.2), ACCENT_GREEN)
for i, item in enumerate(["Normalize flat JSON", "Validate fields", "Remove duplicates"]):
    add_shape(slide, Inches(0.5) + Inches(i * 2.1), Inches(3.3), Inches(1.9), Inches(0.4), BG_CARD,
              item, font_size=9, font_color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

add_arrow_down(slide, Inches(3.3), Inches(3.7), Inches(0.2), ACCENT_GREEN)
add_shape(slide, Inches(1.5), Inches(4.1), Inches(3.5), Inches(0.5), RGBColor(0x1B,0x5E,0x20),
          "S3_KBConfig (~220 flat JSON configs)", font_size=11, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# DynamoDB dual-write for configs
add_shape(slide, Inches(1.5), Inches(4.7), Inches(3.5), Inches(0.4), RGBColor(0x0D,0x47,0xA1),
          "+ DynamoDB StreamingConfigs (dual-write)", font_size=10, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Logs Pipeline (right)
add_rect(slide, Inches(6.8), Inches(0.8), Inches(6.2), Inches(0.45), ACCENT_ORANGE,
         "CloudWatch Metrics Pipeline (every 1h)", font_size=13, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
log_y = Inches(1.5)
add_shape(slide, Inches(7.0), log_y, Inches(1.6), Inches(0.65), RGBColor(0x88,0x0E,0x4F),
          "EventBridge\nScheduler (1h)", font_size=10, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_arrow_right(slide, Inches(8.6), log_y + Inches(0.3), Inches(0.3), ACCENT_YELLOW)
add_shape(slide, Inches(9.0), log_y, Inches(1.6), Inches(0.65), RGBColor(0xE6,0x51,0x00),
          "Pipeline_Logs\nLambda", font_size=10, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_arrow_right(slide, Inches(10.6), log_y + Inches(0.3), Inches(0.3), ACCENT_YELLOW)
add_shape(slide, Inches(11.0), log_y, Inches(1.8), Inches(0.65), RGBColor(0x88,0x0E,0x4F),
          "CloudWatch\nMetrics", font_size=10, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_arrow_down(slide, Inches(9.8), Inches(2.15), Inches(0.2), ACCENT_ORANGE)
for i, (sev, c) in enumerate([("INFO", ACCENT_GREEN), ("WARNING", ACCENT_YELLOW), ("ERROR", ACCENT_ORANGE), ("CRITICAL", ACCENT_RED)]):
    add_rect(slide, Inches(7.0) + Inches(i * 1.55), Inches(2.5), Inches(1.4), Inches(0.4), c,
             sev, font_size=10, font_color=DARK_TEXT if sev in ("WARNING", "INFO") else WHITE, bold=True, align=PP_ALIGN.CENTER)

add_arrow_down(slide, Inches(9.8), Inches(2.9), Inches(0.2), ACCENT_ORANGE)
add_shape(slide, Inches(7.2), Inches(3.3), Inches(5.5), Inches(1.2), BG_CARD)
add_text(slide, Inches(7.4), Inches(3.35), Inches(5), Inches(0.3),
         "Structured Event:", font_size=11, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, Inches(7.4), Inches(3.65), Inches(5), Inches(0.8),
                ["timestamp | channel | severity | error_type",
                 "description | probable_cause",
                 "recommended_fix | source_service"],
                font_size=10, color=LIGHT_GRAY)

add_arrow_down(slide, Inches(9.8), Inches(4.5), Inches(0.2), ACCENT_GREEN)
add_shape(slide, Inches(7.8), Inches(4.9), Inches(4.0), Inches(0.5), RGBColor(0x1B,0x5E,0x20),
          "S3_KBLogs (Structured Events)", font_size=11, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# DynamoDB dual-write for logs
add_shape(slide, Inches(7.8), Inches(5.5), Inches(4.0), Inches(0.4), RGBColor(0x0D,0x47,0xA1),
          "+ DynamoDB StreamingLogs (dual-write, TTL 30d)", font_size=10, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# SNS Alerts
add_shape(slide, Inches(7.8), Inches(6.0), Inches(4.0), Inches(0.4), RGBColor(0x88,0x0E,0x4F),
          "+ SNS Proactive Alerts (ERROR/CRITICAL)", font_size=10, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.5), Inches(6.6), Inches(12), Inches(0.8),
         "Dual-write feeds both Bedrock KB (RAG via S3) and DynamoDB (fast queries). SNS alerts for ERROR/CRITICAL events.",
         font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 8 — DIAGRAM: MULTI-REGION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BLUE)
add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
         "Multi-Region Architecture", font_size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# us-east-1
add_rect(slide, Inches(0.5), Inches(0.9), Inches(5.8), Inches(6.2), RGBColor(0x1A,0x2A,0x4A))
add_rect(slide, Inches(0.5), Inches(0.9), Inches(5.8), Inches(0.5), ACCENT_BLUE,
         "us-east-1  (API / AI / Auth)", font_size=14, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
us_services = [
    ("Amazon Cognito", "User Pool, JWT Auth", RGBColor(0x88,0x0E,0x4F)),
    ("API Gateway REST", "REST endpoints", RGBColor(0xE6,0x51,0x00)),
    ("Lambda Function URL", "Streaming, 5min timeout", RGBColor(0xE6,0x51,0x00)),
    ("Orchestrator Lambda", "Entry point", RGBColor(0xE6,0x51,0x00)),
    ("Bedrock Agent (Claude)", "RAG, Portuguese", RGBColor(0x4A,0x14,0x8C)),
    ("Knowledge Bases", "KB_CONFIG + KB_LOGS", RGBColor(0x1B,0x5E,0x20)),
    ("DynamoDB Tables", "StreamingConfigs + StreamingLogs", RGBColor(0x0D,0x47,0xA1)),
    ("Configurator Lambda", "Create channels + rollback", RGBColor(0xE6,0x51,0x00)),
    ("Exporter Lambda", "DynamoDB queries + S3 fallback", RGBColor(0xE6,0x51,0x00)),
    ("S3 Buckets (5)", "Config, Logs, Audit, Exports, FE", RGBColor(0x1B,0x5E,0x20)),
    ("SNS + EventBridge", "Proactive alerts + Schedulers", RGBColor(0x88,0x0E,0x4F)),
]
for i, (name, desc, color) in enumerate(us_services):
    y = Inches(1.6) + Inches(i * 0.52)
    add_shape(slide, Inches(0.8), y, Inches(2.5), Inches(0.42), color,
              name, font_size=10, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(3.4), y + Inches(0.05), Inches(2.5), Inches(0.35),
             desc, font_size=9, color=LIGHT_GRAY)

# sa-east-1
add_rect(slide, Inches(7.0), Inches(0.9), Inches(5.8), Inches(6.2), RGBColor(0x4A,0x1A,0x1A))
add_rect(slide, Inches(7.0), Inches(0.9), Inches(5.8), Inches(0.5), ACCENT_RED,
         "sa-east-1  (Streaming Services)", font_size=14, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
sa_services = [
    ("AWS MediaLive", "Live channels, RTMP/RTP inputs\nAutomatic failover", RGBColor(0xE6,0x51,0x00)),
    ("AWS MediaPackage V2", "Channel Groups, Channels\nHLS/DASH endpoints", RGBColor(0xE6,0x51,0x00)),
    ("AWS MediaTailor", "Ad insertion\nContent personalization", RGBColor(0xE6,0x51,0x00)),
    ("Amazon CloudFront", "CDN distributions\nContent delivery", RGBColor(0x4A,0x14,0x8C)),
    ("Amazon CloudWatch", "Health metrics\nfrom all 4 services", RGBColor(0x88,0x0E,0x4F)),
]
for i, (name, desc, color) in enumerate(sa_services):
    y = Inches(1.7) + Inches(i * 1.05)
    add_shape(slide, Inches(7.3), y, Inches(2.5), Inches(0.85), color,
              name, font_size=11, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(10.0), y + Inches(0.1), Inches(2.5), Inches(0.7),
             desc, font_size=10, color=LIGHT_GRAY)

add_rect(slide, Inches(6.3), Inches(3.5), Inches(0.7), Inches(0.04), ACCENT_YELLOW)
add_rect(slide, Inches(6.3), Inches(3.7), Inches(0.7), Inches(0.04), ACCENT_YELLOW)
add_text(slide, Inches(6.15), Inches(3.1), Inches(1.0), Inches(0.35),
         "APIs +\nMetrics", font_size=9, color=ACCENT_YELLOW, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 9 — TECHNOLOGY STACK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_PURPLE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
         "Technology Stack", font_size=32, color=WHITE, bold=True)

tech_items = [
    ("Frontend", "HTML/JS, Amazon S3, CloudFront"),
    ("Authentication", "Amazon Cognito (User Pool, admin-only)"),
    ("API", "API Gateway REST + Lambda Function URL"),
    ("AI / ML", "Amazon Bedrock (Claude), RAG with Knowledge Bases"),
    ("Compute", "AWS Lambda (Python 3.9+)"),
    ("Storage", "Amazon S3 (5 buckets) + DynamoDB (fast queries)"),
    ("Streaming", "MediaLive, MediaPackage V2, MediaTailor"),
    ("CDN", "Amazon CloudFront (distributions)"),
    ("Scheduling", "Amazon EventBridge Scheduler"),
    ("Notifications", "Amazon SNS (proactive alerts)"),
    ("IaC", "AWS CDK (Python) — 10 modular stacks"),
    ("Monitoring", "Amazon CloudWatch Metrics"),
]
col1 = tech_items[:6]
col2 = tech_items[6:]
for i, (cat, tech) in enumerate(col1):
    y = Inches(1.3) + Inches(i * 0.95)
    add_shape(slide, Inches(0.8), y, Inches(5.5), Inches(0.75), BG_CARD)
    add_text(slide, Inches(1.0), y + Inches(0.05), Inches(5), Inches(0.3),
             cat, font_size=14, color=ACCENT_PURPLE, bold=True)
    add_text(slide, Inches(1.0), y + Inches(0.38), Inches(5), Inches(0.3),
             tech, font_size=13, color=LIGHT_GRAY)
for i, (cat, tech) in enumerate(col2):
    y = Inches(1.3) + Inches(i * 0.95)
    add_shape(slide, Inches(7.0), y, Inches(5.5), Inches(0.75), BG_CARD)
    add_text(slide, Inches(7.2), y + Inches(0.05), Inches(5), Inches(0.3),
             cat, font_size=14, color=ACCENT_PURPLE, bold=True)
    add_text(slide, Inches(7.2), y + Inches(0.38), Inches(5), Inches(0.3),
             tech, font_size=13, color=LIGHT_GRAY)


# ============================================================
# SLIDE 10 — CHATBOT FEATURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BLUE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
         "Feature 1: RAG-Powered Chatbot", font_size=30, color=WHITE, bold=True)

add_shape(slide, Inches(0.8), Inches(1.3), Inches(6), Inches(5.5), BG_CARD)
add_text(slide, Inches(1.1), Inches(1.5), Inches(5.5), Inches(0.4),
         "How it works", font_size=20, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(1.1), Inches(2.1), Inches(5.5), Inches(3.5),
                ["Operator asks questions in natural language (PT-BR)",
                 "Bedrock Agent queries 2 Knowledge Bases via RAG",
                 "Exporter uses DynamoDB for sub-second data queries",
                 "KB_CONFIG: ~220 channel configurations (flat JSON)",
                 "KB_LOGS: Structured events from CloudWatch metrics",
                 "45+ categorized suggestions in the sidebar",
                 "Health check, comparison, audit history features"],
                font_size=14, color=LIGHT_GRAY)

add_shape(slide, Inches(7.2), Inches(1.3), Inches(5.5), Inches(5.5), BG_CARD)
add_text(slide, Inches(7.5), Inches(1.5), Inches(5), Inches(0.4),
         "Example queries", font_size=20, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(7.5), Inches(2.1), Inches(5), Inches(4),
                ['"Which channels have 1080p resolution?"',
                 '"Show critical errors from the last 24h"',
                 '"What is the configuration of channel X?"',
                 '"Which channels use H.265 codec?"',
                 '"Export all MediaLive channels as CSV"',
                 '"Which channels have active failover?"',
                 '"Summarize issues for channel Y"'],
                font_size=14, color=LIGHT_GRAY)


# ============================================================
# SLIDE 11 — KNOWLEDGE BASE: STREAMING DOCUMENTATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_GREEN)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
         "Knowledge Base: Streaming Documentation & Best Practices", font_size=28, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.5),
         "Beyond channel configs and metrics, the KB also ingests 56 reference documents from Unified Streaming and AWS,\nenabling the chatbot to answer general streaming questions and provide expert guidance.",
         font_size=14, color=LIGHT_GRAY)

# Left column — Unified Streaming docs
add_shape(slide, Inches(0.5), Inches(1.9), Inches(6.0), Inches(5.3), BG_CARD)
add_rect(slide, Inches(0.5), Inches(1.9), Inches(6.0), Inches(0.5), ACCENT_PURPLE,
         "Unified Streaming Documentation (40+ PDFs)", font_size=14, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

us_topics = [
    ("Streaming Protocols", "ABR Streaming, MPEG-DASH, HLS (TS & fMP4),\nHTTP Smooth Streaming (HSS), HTTP Dynamic Streaming (HDS)"),
    ("Content Protection", "Common Encryption (CENC), Multi-DRM (CENC + PIFF),\nCPIX, DRM with Key Rotation, Trans DRM, Watermarking"),
    ("Packaging & Origin", "Packaging for DASH / HLS / HSS / HDS,\nUnified Origin, Dynamic Manifests, Player URLs"),
    ("Advanced Features", "Low Latency (DVB-DASH), Trick Play, Subtitles,\nID3 Tags, Alternate Audio, Tiled Thumbnails"),
    ("Architecture", "Media Processing Overview, Pipeline Configuration,\nPlugin Library, Encoding Services (x264)"),
]
for i, (topic, details) in enumerate(us_topics):
    y = Inches(2.6) + Inches(i * 0.9)
    add_rect(slide, Inches(0.7), y, Inches(0.06), Inches(0.75), ACCENT_PURPLE)
    add_text(slide, Inches(0.95), y, Inches(5.3), Inches(0.3),
             topic, font_size=13, color=ACCENT_PURPLE, bold=True)
    add_text(slide, Inches(0.95), y + Inches(0.3), Inches(5.3), Inches(0.45),
             details, font_size=11, color=LIGHT_GRAY)

# Right column — AWS docs
add_shape(slide, Inches(6.8), Inches(1.9), Inches(6.0), Inches(3.0), BG_CARD)
add_rect(slide, Inches(6.8), Inches(1.9), Inches(6.0), Inches(0.5), ACCENT_ORANGE,
         "AWS Service Documentation (6 PDFs)", font_size=14, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

aws_docs = [
    ("MediaLive", "User Guide + API Reference"),
    ("MediaPackage", "User Guide + API Reference"),
    ("MediaTailor", "API Reference"),
    ("CloudFront", "Developer Guide"),
    ("SPEKE", "Secure Packager & Encoder Key Exchange"),
]
for i, (svc, desc) in enumerate(aws_docs):
    y = Inches(2.6) + Inches(i * 0.45)
    add_rect(slide, Inches(7.0), y, Inches(0.06), Inches(0.35), ACCENT_ORANGE)
    add_text(slide, Inches(7.25), y, Inches(2.2), Inches(0.3),
             svc, font_size=13, color=ACCENT_ORANGE, bold=True)
    add_text(slide, Inches(9.5), y, Inches(3), Inches(0.3),
             desc, font_size=12, color=LIGHT_GRAY)

# Example questions box
add_shape(slide, Inches(6.8), Inches(5.1), Inches(6.0), Inches(2.1), BG_CARD)
add_rect(slide, Inches(6.8), Inches(5.1), Inches(6.0), Inches(0.45), ACCENT_BLUE,
         "Example General Streaming Questions", font_size=13, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_bullet_list(slide, Inches(7.0), Inches(5.65), Inches(5.6), Inches(1.5),
                ['"What is the difference between HLS and DASH?"',
                 '"How does Common Encryption (CENC) work?"',
                 '"Explain ABR streaming and how bitrate adaptation works"',
                 '"What are best practices for low-latency live streaming?"',
                 '"How do I configure DRM with multiple keys?"'],
                font_size=11, color=LIGHT_GRAY)


# ============================================================
# SLIDE 12 — DATA FLOWS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_PURPLE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
         "Main Data Flows", font_size=32, color=WHITE, bold=True)

flows = [
    ("1. Query Flow", ACCENT_BLUE,
     "User  >  Frontend (CloudFront + S3)  >  Lambda Function URL (5min)  >  Bedrock Agent (Claude)  >  Knowledge Bases (RAG)  >  Contextualized Response"),
    ("2. Export Flow", ACCENT_GREEN,
     "User requests export  >  Direct bypass (no Bedrock)  >  Exporter Lambda  >  DynamoDB query (<500ms) with S3 fallback  >  CSV/JSON  >  Pre-signed URL"),
    ("3. Creation Flow", ACCENT_ORANGE,
     "User requests channel  >  Bedrock collects parameters  >  Configurator Lambda  >  Template cloning (MPV2 endpoints)  >  4 orchestrated steps  >  Rollback on failure"),
]
for i, (title, color, desc) in enumerate(flows):
    y = Inches(1.3) + Inches(i * 2.0)
    add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(1.7), BG_CARD)
    add_rect(slide, Inches(0.5), y, Inches(0.08), Inches(1.7), color)
    add_text(slide, Inches(0.8), y + Inches(0.1), Inches(5), Inches(0.35),
             title, font_size=18, color=color, bold=True)
    add_text(slide, Inches(0.8), y + Inches(0.6), Inches(11.5), Inches(0.9),
             desc, font_size=14, color=LIGHT_GRAY)


# ============================================================
# SLIDE 12 — CDK INFRASTRUCTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_ORANGE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
         "Infrastructure as Code — AWS CDK (Python)", font_size=28, color=WHITE, bold=True)

stacks = [
    ("main_stack", "Overall stack orchestration"),
    ("api_stack", "API Gateway REST + Orchestrator Lambda"),
    ("bedrock_agent_stack", "Bedrock Agent + Portuguese instructions"),
    ("bedrock_kb_stack", "Knowledge Bases (Config + Logs)"),
    ("s3_stack", "5 S3 buckets with lifecycle policies"),
    ("dynamodb_stack", "DynamoDB tables (StreamingConfigs + Logs)"),
    ("configuradora_stack", "Configurator Lambda (create channels)"),
    ("exportadora_stack", "Exporter Lambda (CSV/JSON)"),
    ("pipeline_config_stack", "Config Pipeline + EventBridge (6h)"),
    ("pipeline_logs_stack", "Logs Pipeline + EventBridge (1h)"),
    ("frontend_stack", "S3 + CloudFront for hosting"),
]
for i, (name, desc) in enumerate(stacks):
    col = 0 if i < 6 else 1
    row = i if i < 6 else i - 6
    x = Inches(0.8) + Inches(col * 6.3)
    y = Inches(1.3) + Inches(row * 0.95)
    add_shape(slide, x, y, Inches(5.8), Inches(0.8), BG_CARD)
    add_rect(slide, x, y, Inches(0.08), Inches(0.8), ACCENT_ORANGE)
    add_text(slide, x + Inches(0.3), y + Inches(0.02), Inches(5), Inches(0.35),
             name, font_size=14, color=ACCENT_ORANGE, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.38), Inches(5), Inches(0.35),
             desc, font_size=12, color=LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
         "Multi-region deploy: us-east-1 (API / Bedrock / Cognito)  +  sa-east-1 (Streaming Services)",
         font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 13 — SECURITY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_RED)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
         "Security & Authentication", font_size=32, color=WHITE, bold=True)

security_items = [
    ("Amazon Cognito", "User Pool with self-signup disabled, admin-only user creation, JWT authentication"),
    ("IAM Roles", "Least privilege for each Lambda — minimum required permissions per table/bucket"),
    ("S3 Lifecycle", "Exports: 24h TTL | Audit: 365-day retention | KBs: versioning enabled"),
    ("DynamoDB Security", "Encryption at rest (AWS managed key), point-in-time recovery on StreamingConfigs"),
    ("CloudFront OAI", "Origin Access Identity for secure S3 Frontend access"),
    ("Audit Trail", "Complete record of all operations stored in S3_Audit"),
]
for i, (title, desc) in enumerate(security_items):
    y = Inches(1.3) + Inches(i * 1.0)
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.85), BG_CARD)
    add_rect(slide, Inches(0.8), y, Inches(0.08), Inches(0.85), ACCENT_RED)
    add_text(slide, Inches(1.2), y + Inches(0.05), Inches(10), Inches(0.3),
             title, font_size=16, color=ACCENT_RED, bold=True)
    add_text(slide, Inches(1.2), y + Inches(0.4), Inches(10), Inches(0.35),
             desc, font_size=14, color=LIGHT_GRAY)


# ============================================================
# SLIDE 14 — COST ESTIMATE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_GREEN)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
         "Monthly Cost Estimate", font_size=32, color=WHITE, bold=True)

costs = [
    ("Amazon Bedrock (Claude)", "$30 — $100", ACCENT_PURPLE),
    ("Amazon S3 (5 buckets)", "$1 — $5", ACCENT_GREEN),
    ("Amazon DynamoDB (2 tables + GSIs)", "~$3", RGBColor(0x0D, 0x47, 0xA1)),
    ("AWS Lambda (5 functions)", "$0 — $5", ACCENT_ORANGE),
    ("Amazon CloudFront", "$1 — $5", ACCENT_BLUE),
    ("SNS + Cognito + API GW + EventBridge", "Minimal (~$2)", LIGHT_GRAY),
]
add_rect(slide, Inches(2), Inches(1.5), Inches(5), Inches(0.6), ACCENT_BLUE,
         "Service", font_size=16, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(7), Inches(1.5), Inches(4), Inches(0.6), ACCENT_BLUE,
         "Estimated Cost / month", font_size=16, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
for i, (service, cost, color) in enumerate(costs):
    y = Inches(2.2) + Inches(i * 0.7)
    add_shape(slide, Inches(2), y, Inches(5), Inches(0.6), BG_CARD)
    add_text(slide, Inches(2.2), y + Inches(0.1), Inches(4.5), Inches(0.4),
             service, font_size=15, color=color)
    add_shape(slide, Inches(7), y, Inches(4), Inches(0.6), BG_CARD)
    add_text(slide, Inches(7.2), y + Inches(0.1), Inches(3.5), Inches(0.4),
             cost, font_size=15, color=WHITE)

y_total = Inches(2.2) + Inches(len(costs) * 0.7) + Inches(0.2)
add_rect(slide, Inches(2), y_total, Inches(5), Inches(0.7), ACCENT_GREEN,
         "ESTIMATED TOTAL", font_size=18, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(7), y_total, Inches(4), Inches(0.7), ACCENT_GREEN,
         "$38 — $125 / month", font_size=18, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(2), y_total + Inches(1.0), Inches(9), Inches(0.5),
         "* Streaming service costs (MediaLive, MediaPackage, etc.) are separate and usage-dependent.",
         font_size=12, color=LIGHT_GRAY)


# ============================================================
# SLIDE 15 — LAMBDA DETAIL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_ORANGE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
         "Lambda Functions — Detail", font_size=30, color=WHITE, bold=True)

lambdas_detail = [
    ("Orchestrator Lambda", "Main entry point. Invokes Bedrock Agent for queries. Direct bypass for exports and config downloads.", ACCENT_ORANGE),
    ("Configurator Lambda", "Creates/modifies AWS resources with template cloning. Channel comparison, audit history, health check dashboard. Orchestrated creation with rollback.", ACCENT_RED),
    ("Exporter Lambda", "Queries DynamoDB (primary) with S3 fallback. Generates CSV/JSON with pre-signed URLs. Sub-second query latency.", ACCENT_GREEN),
    ("Pipeline_Config", "Collects configs from 4 services every 6h. Dual-write: S3 (RAG) + DynamoDB StreamingConfigs (fast queries).", ACCENT_BLUE),
    ("Pipeline_Logs", "Collects CloudWatch metrics every 1h. Classifies severity. Dual-write: S3 + DynamoDB StreamingLogs (TTL 30d). Proactive SNS alerts.", ACCENT_PURPLE),
]
for i, (name, desc, color) in enumerate(lambdas_detail):
    y = Inches(1.2) + Inches(i * 1.2)
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(1.05), BG_CARD)
    add_rect(slide, Inches(0.8), y, Inches(0.08), Inches(1.05), color)
    add_text(slide, Inches(1.2), y + Inches(0.05), Inches(10), Inches(0.35),
             name, font_size=17, color=color, bold=True)
    add_text(slide, Inches(1.2), y + Inches(0.45), Inches(10.5), Inches(0.5),
             desc, font_size=13, color=LIGHT_GRAY)


# ============================================================
# SLIDE 16 — ROADMAP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_YELLOW)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
         "Next Steps / Roadmap", font_size=32, color=WHITE, bold=True)

roadmap = [
    ("Completed", "DynamoDB as fast query layer — sub-second latency (was 30-120s)", ACCENT_GREEN),
    ("Completed", "CloudWatch Metrics ingestion with severity classification", ACCENT_GREEN),
    ("Completed", "Proactive SNS alerts for ERROR/CRITICAL events", ACCENT_GREEN),
    ("Completed", "Channel comparison, audit history, mass health check", ACCENT_GREEN),
    ("Short term", "Visual channel health dashboard in real time", ACCENT_BLUE),
    ("Mid term", "Predictive failure analysis using metrics history", ACCENT_PURPLE),
]
for i, (phase, desc, color) in enumerate(roadmap):
    y = Inches(1.3) + Inches(i * 0.95)
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.8), BG_CARD)
    add_rect(slide, Inches(0.8), y, Inches(0.08), Inches(0.8), color)
    add_rect(slide, Inches(1.1), y + Inches(0.2), Inches(1.8), Inches(0.4), color,
             phase, font_size=12, font_color=DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(3.2), y + Inches(0.2), Inches(9), Inches(0.4),
             desc, font_size=15, color=LIGHT_GRAY)


# ============================================================
# SLIDE 17 — CLOSING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BLUE)
add_rect(slide, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08), ACCENT_PURPLE)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
         "Thank You!", font_size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(2), Inches(3.0), Inches(9), Inches(3.5), BG_CARD)
add_bullet_list(slide, Inches(2.5), Inches(3.3), Inches(8), Inches(2.8),
                ["AI-powered chatbot for NOC simplifying management of ~220 streaming channels",
                 "Natural language queries with RAG over real data",
                 "DynamoDB fast queries (<500ms) replacing S3 scans (30-120s) with dual-write",
                 "Automated channel creation with template cloning and rollback",
                 "Proactive SNS alerts, health check dashboard, channel comparison, audit history",
                 "Serverless, scalable, low-cost architecture ($38-125/month)"],
                font_size=16, color=LIGHT_GRAY)

add_text(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
         "Questions?", font_size=28, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
output_path = "apresentacao/streaming-platform-ai.pptx"
prs.save(output_path)
print(f"PowerPoint saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
